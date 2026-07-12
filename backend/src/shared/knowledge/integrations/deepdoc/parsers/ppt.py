from __future__ import annotations

from io import BytesIO
import logging

# Adapted from RAGFlow deepdoc/parser/ppt_parser.py


class RAGFlowPptParser:
    def __init__(self):
        self._shape_cache = {}

    @staticmethod
    def _import_presentation():
        from pptx import Presentation

        return Presentation

    def _sort_shapes(self, shapes):
        cache_key = id(shapes)
        if cache_key not in self._shape_cache:
            self._shape_cache[cache_key] = sorted(
                shapes,
                key=lambda shape: (
                    (shape.top if shape.top is not None else 0) // 10,
                    shape.left if shape.left is not None else 0,
                ),
            )
        return self._shape_cache[cache_key]

    @staticmethod
    def _get_bulleted_text(paragraph):
        is_bulleted = (
            bool(paragraph._p.xpath("./a:pPr/a:buChar"))
            or bool(paragraph._p.xpath("./a:pPr/a:buAutoNum"))
            or bool(paragraph._p.xpath("./a:pPr/a:buBlip"))
        )
        return f"{'  ' * paragraph.level}.{paragraph.text}" if is_bulleted else paragraph.text

    def _extract(self, shape):
        try:
            if hasattr(shape, "has_text_frame") and shape.has_text_frame:
                text_frame = shape.text_frame
                texts = []
                for paragraph in text_frame.paragraphs:
                    if paragraph.text.strip():
                        texts.append(self._get_bulleted_text(paragraph))
                return "\n".join(texts)

            try:
                shape_type = shape.shape_type
            except NotImplementedError:
                return shape.text.strip() if hasattr(shape, "text") else ""

            if shape_type == 19:
                table = shape.table
                rows = []
                for row_idx in range(1, len(table.rows)):
                    rows.append(
                        "; ".join(
                            f"{table.cell(0, col_idx).text}: {table.cell(row_idx, col_idx).text}"
                            for col_idx in range(len(table.columns))
                            if table.cell(row_idx, col_idx)
                        )
                    )
                return "\n".join(rows)

            if shape_type == 6:
                texts = []
                for child_shape in self._sort_shapes(shape.shapes):
                    text = self._extract(child_shape)
                    if text:
                        texts.append(text)
                return "\n".join(texts)

            return ""
        except Exception as exc:
            logging.error("Error processing PPT shape: %s", exc)
            return ""

    def __call__(self, fnm, from_page=0, to_page=None, callback=None):
        Presentation = self._import_presentation()
        ppt = Presentation(fnm) if isinstance(fnm, str) else Presentation(BytesIO(fnm))
        texts_by_slide = []
        self.total_page = len(ppt.slides)
        upper = self.total_page if to_page is None else to_page
        for slide_index, slide in enumerate(ppt.slides):
            if slide_index < from_page:
                continue
            if slide_index >= upper:
                break
            texts = []
            for shape in self._sort_shapes(slide.shapes):
                text = self._extract(shape)
                if text:
                    texts.append(text)
            texts_by_slide.append("\n".join(texts))
        return texts_by_slide
